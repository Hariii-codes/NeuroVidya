#!/usr/bin/env python3
"""
NeuroVidya PowerPoint Presentation Generator
Creates a dyslexia-friendly, professional 10-slide presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Color Scheme (Dyslexia-Friendly)
COLORS = {
    'primary': RGBColor(59, 130, 246),      # Calm blue
    'secondary': RGBColor(139, 92, 246),    # Soft purple
    'highlight': RGBColor(251, 191, 36),    # Warm yellow
    'background': RGBColor(250, 246, 240),  # Cream
    'text': RGBColor(45, 55, 72),           # Dark gray
    'white': RGBColor(255, 255, 255),
    'success': RGBColor(34, 197, 94),       # Green
    'light_gray': RGBColor(229, 231, 235),
}

def set_slide_background(slide, color):
    """Set solid background color for slide"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, slide_layout):
    """Create Title Slide"""
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['primary'])

    # Clear existing placeholders
    for shape in slide.placeholders:
        element = shape.element
        element.getparent().remove(element)

    # Main Title
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = "NEUROVIDYA"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['white']
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "A Dyslexia-First Adaptive Learning Platform"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = COLORS['white']
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Tagline
    tagline_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.5))
    tagline_frame = tagline_box.text_frame
    tagline_frame.text = '"Transforming Learning Through Accessibility-First Design"'
    tagline_para = tagline_frame.paragraphs[0]
    tagline_para.font.size = Pt(20)
    tagline_para.font.color.rgb = COLORS['highlight']
    tagline_para.alignment = PP_ALIGN.CENTER

    # Team Members
    team_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.5), Inches(7), Inches(1.5))
    team_frame = team_box.text_frame
    team_frame.word_wrap = True
    team_frame.text = "Team Members:\n• Project Lead & Frontend Developer\n• Backend Developer\n• UI/UX Designer\n• Research & Documentation"
    team_para = team_frame.paragraphs[0]
    team_para.font.size = Pt(18)
    team_para.font.color.rgb = COLORS['white']
    for para in team_frame.paragraphs:
        para.font.size = Pt(18)
        para.font.color.rgb = COLORS['white']

    # Institution
    inst_box = slide.shapes.add_textbox(Inches(1), Inches(7.2), Inches(8), Inches(0.4))
    inst_frame = inst_box.text_frame
    inst_frame.text = "Department of Computer Science & Engineering | Academic Year 2025-2026"
    inst_para = inst_frame.paragraphs[0]
    inst_para.font.size = Pt(16)
    inst_para.font.color.rgb = COLORS['white']
    inst_para.alignment = PP_ALIGN.CENTER

    return slide

def add_intro_slide(prs, slide_layout):
    """Create Introduction Slide"""
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['background'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Introduction to NeuroVidya"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']

    # What is NeuroVidya
    content_left = Inches(0.7)
    content_top = Inches(1.8)

    # Section 1
    section1 = slide.shapes.add_textbox(content_left, content_top, Inches(8.6), Inches(1.2))
    f1 = section1.text_frame
    f1.text = "✦ What is NeuroVidya?"
    p1 = f1.paragraphs[0]
    p1.font.size = Pt(24)
    p1.font.bold = True
    p1.font.color.rgb = COLORS['primary']
    p1.space_after = Pt(8)

    f1.text += "\n• Comprehensive web-based adaptive learning platform for dyslexia"
    p2 = f1.paragraphs[1]
    p2.font.size = Pt(20)
    p2.font.color.rgb = COLORS['text']
    p2.space_before = Pt(4)

    f1.text += "\n• Dyslexia-first design — not retrofit accessibility"
    p3 = f1.paragraphs[2]
    p3.font.size = Pt(20)
    p3.font.color.rgb = COLORS['text']

    # Section 2
    section2 = slide.shapes.add_textbox(content_left, Inches(3.2), Inches(8.6), Inches(1.2))
    f2 = section2.text_frame
    f2.text = "✦ Why It Matters"
    p2_1 = f2.paragraphs[0]
    p2_1.font.size = Pt(24)
    p2_1.font.bold = True
    p2_1.font.color.rgb = COLORS['primary']

    f2.text += "\n• Dyslexia affects 10-20% of the global population"
    f2.text += "\n• Traditional platforms designed for neurotypical users"
    f2.text += "\n• Accessibility features often added as afterthoughts"
    for i in [1, 2, 3]:
        f2.paragraphs[i].font.size = Pt(20)
        f2.paragraphs[i].font.color.rgb = COLORS['text']

    # Section 3
    section3 = slide.shapes.add_textbox(content_left, Inches(4.6), Inches(8.6), Inches(1.5))
    f3 = section3.text_frame
    f3.text = "✦ Core Philosophy"
    p3_1 = f3.paragraphs[0]
    p3_1.font.size = Pt(24)
    p3_1.font.bold = True
    p3_1.font.color.rgb = COLORS['primary']

    f3.text += "\n• Accessibility from the ground up, not as an add-on"
    f3.text += "\n• Research-backed design principles"
    f3.text += "\n• Empowering users through extensive customization"
    f3.text += "\n• Privacy-first approach with browser-native features"
    for i in [1, 2, 3, 4]:
        f3.paragraphs[i].font.size = Pt(20)
        f3.paragraphs[i].font.color.rgb = COLORS['text']

    # Impact banner
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(6.5), Inches(8.6), Inches(0.8))
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLORS['highlight']
    banner.line.color.rgb = COLORS['highlight']

    banner_frame = banner.text_frame
    banner_frame.text = "Impact: Democratizing access to quality education"
    banner_frame.word_wrap = True
    banner_para = banner_frame.paragraphs[0]
    banner_para.font.size = Pt(22)
    banner_para.font.bold = True
    banner_para.font.color.rgb = COLORS['text']
    banner_para.alignment = PP_ALIGN.CENTER

    return slide

def add_problem_slide(prs, slide_layout):
    """Create Problem Statement Slide"""
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['background'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Problem Statement"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']

    # Core Problem Banner
    problem_banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(8.6), Inches(1))
    problem_banner.fill.solid()
    problem_banner.fill.fore_color.rgb = RGBColor(239, 68, 68)  # Red
    problem_banner.line.color.rgb = RGBColor(239, 68, 68)

    pb_frame = problem_banner.text_frame
    pb_frame.text = "Educational platforms are designed for the neurotypical majority,"
    pb_frame.text += "\nleaving 10-20% of the population with dyslexia struggling."
    pb_para = pb_frame.paragraphs[0]
    pb_para.font.size = Pt(20)
    pb_para.font.bold = True
    pb_para.font.color.rgb = COLORS['white']
    pb_para.alignment = PP_ALIGN.CENTER
    pb_frame.paragraphs[1].font.size = Pt(20)
    pb_frame.paragraphs[1].font.bold = True
    pb_frame.paragraphs[1].font.color.rgb = COLORS['white']
    pb_frame.paragraphs[1].alignment = PP_ALIGN.CENTER

    # Challenges - Left Column
    left_col = slide.shapes.add_textbox(Inches(0.7), Inches(2.8), Inches(4.1), Inches(2.8))
    lf = left_col.text_frame

    lf.text = "❌ Visual Processing"
    lf.paragraphs[0].font.size = Pt(22)
    lf.paragraphs[0].font.bold = True
    lf.paragraphs[0].font.color.rgb = RGBColor(239, 68, 68)

    lf.text += "\n• Visual crowding & pattern interference"
    lf.text += "\n• Difficulty with left-to-right scanning"
    lf.text += "\n• Eye strain from poor typography"
    for i in [1, 2, 3]:
        lf.paragraphs[i].font.size = Pt(18)
        lf.paragraphs[i].font.color.rgb = COLORS['text']

    # Challenges - Right Column
    right_col = slide.shapes.add_textbox(Inches(5.2), Inches(2.8), Inches(4.1), Inches(2.8))
    rf = right_col.text_frame

    rf.text = "❌ Phonological Processing"
    rf.paragraphs[0].font.size = Pt(22)
    rf.paragraphs[0].font.bold = True
    rf.paragraphs[0].font.color.rgb = RGBColor(239, 68, 68)

    rf.text += "\n• Difficulty breaking words into sounds"
    rf.text += "\n• Challenges with sound-symbol correspondence"
    rf.text += "\n• Poor spelling and decoding abilities"
    for i in [1, 2, 3]:
        rf.paragraphs[i].font.size = Pt(18)
        rf.paragraphs[i].font.color.rgb = COLORS['text']

    # Platform Issues
    platform_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.8), Inches(8.6), Inches(1.2))
    pf = platform_box.text_frame
    pf.text = "Existing Platform Issues:"
    pf.paragraphs[0].font.size = Pt(22)
    pf.paragraphs[0].font.bold = True
    pf.paragraphs[0].font.color.rgb = COLORS['secondary']

    issues = "• Accessibility as add-on • Limited customization • Expensive • No pronunciation feedback"
    pf.text += f"\n{issues}"
    pf.paragraphs[1].font.size = Pt(18)
    pf.paragraphs[1].font.color.rgb = COLORS['text']

    # Research Question
    rq_box = slide.shapes.add_textbox(Inches(0.7), Inches(7.2), Inches(8.6), Inches(0.6))
    rqf = rq_box.text_frame
    rqf.text = 'Research: How can a platform optimize reading comprehension specifically for dyslexic learners?'
    rq_para = rqf.paragraphs[0]
    rq_para.font.size = Pt(18)
    rq_para.font.italic = True
    rq_para.font.color.rgb = COLORS['primary']

    return slide

def add_lit_review_slide(prs, slide_layout):
    """Create Literature Review Slide"""
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['background'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Literature Review & Research Foundation"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']

    # Theoretical Frameworks
    theories = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(2.5))
    tf = theories.text_frame

    tf.text = "Theoretical Frameworks"
    tf.paragraphs[0].font.size = Pt(24)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLORS['primary']

    frameworks = [
        ("Dual Coding Theory", "Paivio, 1971 — Visual + verbal channels enhance learning"),
        ("Cognitive Load Theory", "Sweller, 1988 — Managing visual overwhelm through line focus"),
        ("Phonological Deficit", "Multi-sensory reading support addresses core challenges"),
    ]

    for i, (name, desc) in enumerate(frameworks, 1):
        tf.add_paragraph()
        tf.paragraphs[i].text = f"• {name}: {desc}"
        tf.paragraphs[i].font.size = Pt(18)
        tf.paragraphs[i].font.color.rgb = COLORS['text']
        tf.paragraphs[i].space_before = Pt(4)

    # Comparison Table
    table_left = Inches(0.7)
    table_top = Inches(4.2)

    # Table Header
    header_box = slide.shapes.add_textbox(table_left, table_top, Inches(8.6), Inches(0.5))
    hf = header_box.text_frame
    hf.text = "Existing Solutions vs Limitations"
    hf.paragraphs[0].font.size = Pt(22)
    hf.paragraphs[0].font.bold = True
    hf.paragraphs[0].font.color.rgb = COLORS['secondary']

    # Table content
    table_data = [
        ("Microsoft Immersive Reader", "Good TTS", "Limited customization, requires Office 365"),
        ("Learning Ally", "Extensive library", "Primarily audio, lacks interaction"),
        ("Google Read & Write", "Browser convenience", "Extension limitations"),
        ("Kurzweil 3000", "Established platform", "Expensive, dated interface"),
    ]

    table_top2 = Inches(4.8)
    row_height = Inches(0.45)

    for i, (solution, strength, limitation) in enumerate(table_data):
        row_box = slide.shapes.add_textbox(table_left, table_top2 + (i * row_height), Inches(8.6), row_height)
        rf = row_box.text_frame

        rf.text = f"{solution}"
        rf.paragraphs[0].font.size = Pt(16)
        rf.paragraphs[0].font.bold = True
        rf.paragraphs[0].font.color.rgb = COLORS['primary']

        rf.text += f" | ✓ {strength} | ✗ {limitation}"
        rf.paragraphs[0].font.size = Pt(14)
        rf.paragraphs[0].font.color.rgb = COLORS['text']

    # Research Gap
    gap_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(6.8), Inches(8.6), Inches(0.8))
    gap_box.fill.solid()
    gap_box.fill.fore_color.rgb = COLORS['secondary']
    gap_box.line.color.rgb = COLORS['secondary']

    gap_frame = gap_box.text_frame
    gap_frame.text = "Research Gap: No comprehensive, open-source, dyslexia-first platform with AI integration"
    gap_para = gap_frame.paragraphs[0]
    gap_para.font.size = Pt(18)
    gap_para.font.bold = True
    gap_para.font.color.rgb = COLORS['white']
    gap_para.alignment = PP_ALIGN.CENTER

    return slide

def add_objectives_slide(prs, slide_layout):
    """Create Objectives Slide"""
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['background'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Project Objectives"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']

    # Primary Objectives
    obj_left = Inches(0.7)
    obj_top = Inches(1.6)

    objectives = [
        ("🎯 Dyslexia-First Design System", "Comprehensive component library with research-backed typography"),
        ("🎯 Advanced Phonetic Support", "Syllable chunking, AI phonetic breakdown, real-time feedback"),
        ("🎯 Multi-Modal Learning", "TTS, visual diagrams, interactive games, progress tracking"),
        ("🎯 Accessibility & Usability", "100% keyboard navigation, WCAG AA compliance"),
    ]

    for i, (title, desc) in enumerate(objectives):
        obj_box = slide.shapes.add_textbox(obj_left, obj_top + (i * Inches(1.1)), Inches(8.6), Inches(1))
        of = obj_box.text_frame

        of.text = title
        of.paragraphs[0].font.size = Pt(22)
        of.paragraphs[0].font.bold = True
        of.paragraphs[0].font.color.rgb = COLORS['primary']

        of.text += f"\n{desc}"
        of.paragraphs[1].font.size = Pt(18)
        of.paragraphs[1].font.color.rgb = COLORS['text']

    # Success Metrics
    metrics_box = slide.shapes.add_textbox(obj_left, Inches(6.2), Inches(8.6), Inches(0.4))
    mf = metrics_box.text_frame
    mf.text = "Success Metrics Achieved:"
    mf.paragraphs[0].font.size = Pt(22)
    mf.paragraphs[0].font.bold = True
    mf.paragraphs[0].font.color.rgb = COLORS['success']

    metrics = [
        ("WCAG 2.1 AA Compliance", "✅ 100%"),
        ("Test Coverage", "✅ 79%"),
        ("Page Load Time", "✅ 2.8s"),
        ("Design Components", "✅ 20+"),
    ]

    metrics_top = Inches(6.7)
    for i, (metric, status) in enumerate(metrics):
        m_box = slide.shapes.add_textbox(obj_left + (i * Inches(2.2)), metrics_top, Inches(2), Inches(0.8))
        m_frame = m_box.text_frame
        m_frame.text = f"{metric}\n{status}"
        m_frame.paragraphs[0].font.size = Pt(16)
        m_frame.paragraphs[0].font.bold = True
        m_frame.paragraphs[0].font.color.rgb = COLORS['text']
        m_frame.paragraphs[1].font.size = Pt(18)
        m_frame.paragraphs[1].font.color.rgb = COLORS['success']
        m_frame.paragraphs[1].alignment = PP_ALIGN.CENTER

    return slide

def add_methodology_slide(prs, slide_layout):
    """Create Methodology & Requirements Slide"""
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['background'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Methodology & Software Requirements"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']

    # Development Approach
    approach_left = Inches(0.7)

    approach = slide.shapes.add_textbox(approach_left, Inches(1.5), Inches(8.6), Inches(0.6))
    af = approach.text_frame
    af.text = "Agile Development with 2-Week Sprints | Component-First Architecture"
    af.paragraphs[0].font.size = Pt(22)
    af.paragraphs[0].font.bold = True
    af.paragraphs[0].font.color.rgb = COLORS['secondary']
    af.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Development Phases
    phases = [
        ("Weeks 1-2", "Research & Planning", "✅"),
        ("Weeks 3-4", "Design System Foundation", "✅"),
        ("Weeks 5-7", "Advanced Features", "✅"),
        ("Weeks 8-10", "Page Development", "✅"),
        ("Weeks 11-12", "Backend Integration", "✅"),
        ("Weeks 13-14", "Testing & QA", "✅"),
        ("Weeks 15-16", "User Validation", "🔄"),
    ]

    phases_top = Inches(2.3)
    for i, (week, activity, status) in enumerate(phases):
        phase_box = slide.shapes.add_textbox(
            approach_left + ((i % 4) * Inches(2.2)),
            phases_top + ((i // 4) * Inches(1.2)),
            Inches(2.1),
            Inches(1)
        )
        pf = phase_box.text_frame

        pf.text = f"{week}"
        pf.paragraphs[0].font.size = Pt(14)
        pf.paragraphs[0].font.bold = True
        pf.paragraphs[0].font.color.rgb = COLORS['primary']

        pf.text += f"\n{activity}"
        pf.paragraphs[1].font.size = Pt(16)
        pf.paragraphs[1].font.color.rgb = COLORS['text']

        pf.text += f"\n{status}"
        pf.paragraphs[2].font.size = Pt(18)
        pf.paragraphs[2].font.alignment = PP_ALIGN.CENTER

    # Requirements Section
    req_top = Inches(5.2)

    # Frontend
    fe_box = slide.shapes.add_textbox(approach_left, req_top, Inches(4.1), Inches(1.5))
    fef = fe_box.text_frame
    fef.text = "Frontend Stack"
    fef.paragraphs[0].font.size = Pt(20)
    fef.paragraphs[0].font.bold = True
    fef.paragraphs[0].font.color.rgb = COLORS['primary']

    fe_tech = "React 18.2 | TypeScript 5.3 | Vite 5.0\nZustand | TanStack Query | Tesseract.js"
    fef.text += f"\n{fe_tech}"
    for i in [1, 2]:
        fef.paragraphs[i].font.size = Pt(16)
        fef.paragraphs[i].font.color.rgb = COLORS['text']

    # Backend
    be_box = slide.shapes.add_textbox(Inches(5.2), req_top, Inches(4.1), Inches(1.5))
    bef = be_box.text_frame
    bef.text = "Backend Stack"
    bef.paragraphs[0].font.size = Pt(20)
    bef.paragraphs[0].font.bold = True
    bef.paragraphs[0].font.color.rgb = COLORS['primary']

    be_tech = "FastAPI | PostgreSQL | Python 3.9+\nOpenAI GPT-4o-mini | Anthropic Claude"
    bef.text += f"\n{be_tech}"
    for i in [1, 2]:
        bef.paragraphs[i].font.size = Pt(16)
        bef.paragraphs[i].font.color.rgb = COLORS['text']

    # Platform Support
    platform_box = slide.shapes.add_textbox(approach_left, Inches(6.9), Inches(8.6), Inches(0.5))
    plf = platform_box.text_frame
    plf.text = "Platform Support: Windows ✅ | macOS ✅ | Linux ✅ | Android 🔄 | iOS 🔄"
    plf.paragraphs[0].font.size = Pt(18)
    plf.paragraphs[0].font.color.rgb = COLORS['text']
    plf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def add_implementation_slide(prs, slide_layout):
    """Create Implementation Slide"""
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['background'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Implementation Details"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']

    # System Architecture Diagram (Text-based)
    arch_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(8.6), Inches(1.8))
    arch_box.fill.solid()
    arch_box.fill.fore_color.rgb = RGBColor(237, 242, 255)
    arch_box.line.color.rgb = COLORS['primary']

    arch_frame = arch_box.text_frame
    arch_frame.text = "SYSTEM ARCHITECTURE"
    arch_frame.paragraphs[0].font.size = Pt(18)
    arch_frame.paragraphs[0].font.bold = True
    arch_frame.paragraphs[0].font.color.rgb = COLORS['primary']
    arch_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    arch_frame.paragraphs[0].space_after = Pt(8)

    arch_diagram = """
┌─────────────────────────────────────────┐
│  CLIENT LAYER (React + Web Speech API)  │
└─────────────────────────────────────────┘
                    ↕ REST API
┌─────────────────────────────────────────┐
│         API GATEWAY (FastAPI)           │
└─────────────────────────────────────────┘
                    ↕
┌─────────────────────────────────────────┐
│      DATA LAYER (PostgreSQL + Mongo)    │
└─────────────────────────────────────────┘"""

    arch_frame.add_paragraph()
    arch_frame.paragraphs[1].text = arch_diagram
    arch_frame.paragraphs[1].font.size = Pt(14)
    arch_frame.paragraphs[1].font.name = "Courier New"
    arch_frame.paragraphs[1].font.color.rgb = COLORS['text']
    arch_frame.paragraphs[1].alignment = PP_ALIGN.CENTER

    # Key Components
    comp_left = Inches(0.7)
    comp_top = Inches(3.6)

    comp_box = slide.shapes.add_textbox(comp_left, comp_top, Inches(8.6), Inches(0.5))
    cf = comp_box.text_frame
    cf.text = "Key Components Implemented (20+)"
    cf.paragraphs[0].font.size = Pt(22)
    cf.paragraphs[0].font.bold = True
    cf.paragraphs[0].font.color.rgb = COLORS['secondary']

    components = [
        "DyslexiaProvider", "DyslexiaTheme", "DyslexiaButton",
        "ReadingGuide", "PhoneticHighlighter", "PronunciationAnalyzer",
        "QuickSettingsToolbar", "+ 13 more specialized components"
    ]

    comp_list_top = Inches(4.2)
    for i, comp in enumerate(components):
        c_box = slide.shapes.add_textbox(
            comp_left + ((i % 4) * Inches(2.2)),
            comp_list_top + ((i // 4) * Inches(0.5)),
            Inches(2.1),
            Inches(0.5)
        )
        c_frame = c_box.text_frame
        c_frame.text = f"• {comp}"
        c_frame.paragraphs[0].font.size = Pt(16)
        c_frame.paragraphs[0].font.color.rgb = COLORS['text']

    # Key Algorithms
    algo_top = Inches(5.5)
    algo_box = slide.shapes.add_textbox(comp_left, algo_top, Inches(8.6), Inches(0.5))
    alf = algo_box.text_frame
    alf.text = "Key Algorithms"
    alf.paragraphs[0].font.size = Pt(22)
    alf.paragraphs[0].font.bold = True
    alf.paragraphs[0].font.color.rgb = COLORS['secondary']

    # Algorithm 1
    algo1 = slide.shapes.add_textbox(comp_left, Inches(6.1), Inches(4.1), Inches(0.8))
    a1f = algo1.text_frame
    a1f.text = "1. Phonetic Chunking"
    a1f.paragraphs[0].font.size = Pt(18)
    a1f.paragraphs[0].font.bold = True
    a1f.paragraphs[0].font.color.rgb = COLORS['primary']
    a1f.text += "\nHybrid: Syllable rules + AI enhancement"
    a1f.paragraphs[1].font.size = Pt(16)
    a1f.paragraphs[1].font.color.rgb = COLORS['text']

    # Algorithm 2
    algo2 = slide.shapes.add_textbox(Inches(5.2), Inches(6.1), Inches(4.1), Inches(0.8))
    a2f = algo2.text_frame
    a2f.text = "2. Pronunciation Analysis"
    a2f.paragraphs[0].font.size = Pt(18)
    a2f.paragraphs[0].font.bold = True
    a2f.paragraphs[0].font.color.rgb = COLORS['primary']
    a2f.text += "\nBrowser-native speech recognition + Levenshtein distance"
    a2f.paragraphs[1].font.size = Pt(16)
    a2f.paragraphs[1].font.color.rgb = COLORS['text']

    # Performance Metrics
    perf_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, comp_left, Inches(7.1), Inches(8.6), Inches(0.6))
    perf_box.fill.solid()
    perf_box.fill.fore_color.rgb = COLORS['success']
    perf_box.line.color.rgb = COLORS['success']

    perf_frame = perf_box.text_frame
    perf_frame.text = "Performance: 92 Lighthouse Score | 100 Accessibility | 2.8s Load Time | 79% Test Coverage"
    perf_para = perf_frame.paragraphs[0]
    perf_para.font.size = Pt(16)
    perf_para.font.bold = True
    perf_para.font.color.rgb = COLORS['white']
    perf_para.alignment = PP_ALIGN.CENTER

    return slide

def add_application_slide(prs, slide_layout):
    """Create Application & Future Scope Slide"""
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['background'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Application & Future Scope"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']

    # Current Applications
    app_left = Inches(0.7)

    apps_box = slide.shapes.add_textbox(app_left, Inches(1.5), Inches(8.6), Inches(0.5))
    abf = apps_box.text_frame
    abf.text = "Current Applications"
    abf.paragraphs[0].font.size = Pt(24)
    abf.paragraphs[0].font.bold = True
    abf.paragraphs[0].font.color.rgb = COLORS['secondary']

    app_areas = [
        ("🎓 K-12 Schools", "Special education, IEPs, resource rooms"),
        ("🎓 Higher Education", "Disability services, test accommodations"),
        ("💼 Workplace", "Employee training, ADA compliance"),
        ("👤 Individuals", "Self-directed learning, skill development"),
    ]

    app_top = Inches(2.1)
    for i, (area, desc) in enumerate(app_areas):
        app_box = slide.shapes.add_textbox(
            app_left + ((i % 2) * Inches(4.3)),
            app_top + ((i // 2) * Inches(0.7)),
            Inches(4.2),
            Inches(0.7)
        )
        af = app_box.text_frame
        af.text = area
        af.paragraphs[0].font.size = Pt(18)
        af.paragraphs[0].font.bold = True
        af.paragraphs[0].font.color.rgb = COLORS['primary']

        af.text += f"\n{desc}"
        af.paragraphs[1].font.size = Pt(16)
        af.paragraphs[1].font.color.rgb = COLORS['text']

    # Future Enhancements
    future_top = Inches(4.0)
    future_box = slide.shapes.add_textbox(app_left, future_top, Inches(8.6), Inches(0.5))
    fbf = future_box.text_frame
    fbf.text = "Future Enhancements"
    fbf.paragraphs[0].font.size = Pt(24)
    fbf.paragraphs[0].font.bold = True
    fbf.paragraphs[0].font.color.rgb = COLORS['secondary']

    # Short-term
    st_box = slide.shapes.add_textbox(app_left, Inches(4.6), Inches(4.1), Inches(1.2))
    stf = st_box.text_frame
    stf.text = "🔜 Short-term (6-12 months)"
    stf.paragraphs[0].font.size = Pt(20)
    stf.paragraphs[0].font.bold = True
    stf.paragraphs[0].font.color.rgb = COLORS['highlight']

    st_items = "• Mobile apps (iOS/Android)\n• Offline mode support\n• Advanced analytics\n• Multi-language modules"
    stf.text += f"\n{st_items}"
    for i in [1, 2, 3, 4]:
        stf.paragraphs[i].font.size = Pt(16)
        stf.paragraphs[i].font.color.rgb = COLORS['text']

    # Long-term
    lt_box = slide.shapes.add_textbox(Inches(5.2), Inches(4.6), Inches(4.1), Inches(1.2))
    ltf = lt_box.text_frame
    ltf.text = "🚀 Long-term (1-3 years)"
    ltf.paragraphs[0].font.size = Pt(20)
    ltf.paragraphs[0].font.bold = True
    ltf.paragraphs[0].font.color.rgb = COLORS['highlight']

    lt_items = "• Research platform integration\n• AI learning paths\n• LMS integrations\n• Publishing partnerships"
    ltf.text += f"\n{lt_items}"
    for i in [1, 2, 3, 4]:
        ltf.paragraphs[i].font.size = Pt(16)
        ltf.paragraphs[i].font.color.rgb = COLORS['text']

    # Impact Banner
    impact_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, app_left, Inches(6.1), Inches(8.6), Inches(0.8))
    impact_box.fill.solid()
    impact_box.fill.fore_color.rgb = COLORS['primary']
    impact_box.line.color.rgb = COLORS['primary']

    impact_frame = impact_box.text_frame
    impact_frame.text = "Global Impact Potential: 700+ million people with dyslexia worldwide"
    impact_para = impact_frame.paragraphs[0]
    impact_para.font.size = Pt(22)
    impact_para.font.bold = True
    impact_para.font.color.rgb = COLORS['white']
    impact_para.alignment = PP_ALIGN.CENTER

    # Deployment Options
    deploy_box = slide.shapes.add_textbox(app_left, Inches(7.1), Inches(8.6), Inches(0.5))
    dlf = deploy_box.text_frame
    dlf.text = "Deployment: Vercel ✅ | Netlify ✅ | AWS ✅ | Docker Self-Hosted ✅"
    dlf.paragraphs[0].font.size = Pt(18)
    dlf.paragraphs[0].font.color.rgb = COLORS['text']
    dlf.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def add_conclusion_slide(prs, slide_layout):
    """Create Conclusion Slide"""
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['background'])

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Conclusion & Future Directions"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(40)
    title_para.font.bold = True
    title_para.font.color.rgb = COLORS['primary']

    # Paradigm Shift Banner
    shift_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.5), Inches(8.6), Inches(0.7))
    shift_box.fill.solid()
    shift_box.fill.fore_color.rgb = COLORS['secondary']
    shift_box.line.color.rgb = COLORS['secondary']

    shift_frame = shift_box.text_frame
    shift_frame.text = "From Accessibility-as-Accommodation to Accessibility-as-Foundation"
    shift_para = shift_frame.paragraphs[0]
    shift_para.font.size = Pt(22)
    shift_para.font.bold = True
    shift_para.font.color.rgb = COLORS['white']
    shift_para.alignment = PP_ALIGN.CENTER

    # Key Achievements
    achieve_left = Inches(0.7)

    achievements = [
        ("🎯 Technical Excellence", "79% test coverage | WCAG 2.1 AA compliance | Production-ready"),
        ("🎯 Pedagogical Innovation", "Research-backed design | Multi-modal learning | Real-time feedback"),
        ("🎯 User Empowerment", "Extensive customization | Privacy-first | Open-source"),
    ]

    achieve_top = Inches(2.5)
    for i, (title, desc) in enumerate(achievements):
        ach_box = slide.shapes.add_textbox(achieve_left, achieve_top + (i * Inches(0.9)), Inches(8.6), Inches(0.9))
        af = ach_box.text_frame

        af.text = title
        af.paragraphs[0].font.size = Pt(20)
        af.paragraphs[0].font.bold = True
        af.paragraphs[0].font.color.rgb = COLORS['primary']

        af.text += f"\n{desc}"
        af.paragraphs[1].font.size = Pt(18)
        af.paragraphs[1].font.color.rgb = COLORS['text']

    # Research Contribution
    research_top = Inches(5.4)
    research_box = slide.shapes.add_textbox(achieve_left, research_top, Inches(8.6), Inches(0.5))
    rf = research_box.text_frame
    rf.text = "Research Contributions"
    rf.paragraphs[0].font.size = Pt(22)
    rf.paragraphs[0].font.bold = True
    rf.paragraphs[0].font.color.rgb = COLORS['secondary']

    contributions = [
        "First comprehensive dyslexia-first design system",
        "Hybrid phonetic chunking algorithm",
        "Browser-native pronunciation analysis",
        "Open-source accessibility template",
    ]

    contrib_top = Inches(6.0)
    for i, contrib in enumerate(contributions):
        c_box = slide.shapes.add_textbox(
            achieve_left + ((i % 2) * Inches(4.3)),
            contrib_top + ((i // 2) * Inches(0.5)),
            Inches(4.2),
            Inches(0.5)
        )
        cf = c_box.text_frame
        cf.text = f"• {contrib}"
        cf.paragraphs[0].font.size = Pt(16)
        cf.paragraphs[0].font.color.rgb = COLORS['text']

    # Call to Action
    cta_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, achieve_left, Inches(7.1), Inches(8.6), Inches(0.6))
    cta_box.fill.solid()
    cta_box.fill.fore_color.rgb = COLORS['highlight']
    cta_box.line.color.rgb = COLORS['highlight']

    cta_frame = cta_box.text_frame
    cta_frame.text = "Join us in making education accessible to all | www.neurovidya.com"
    cta_para = cta_frame.paragraphs[0]
    cta_para.font.size = Pt(20)
    cta_para.font.bold = True
    cta_para.font.color.rgb = COLORS['text']
    cta_para.alignment = PP_ALIGN.CENTER

    return slide

def add_thankyou_slide(prs, slide_layout):
    """Create Thank You / Q&A Slide"""
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide, COLORS['primary'])

    # Clear existing placeholders
    for shape in slide.placeholders:
        element = shape.element
        element.getparent().remove(element)

    # Thank You
    thanks_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(0.8))
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "Thank You!"
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.font.size = Pt(54)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = COLORS['white']
    thanks_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(0.6))
    sub_frame = sub_box.text_frame
    sub_frame.text = "Questions & Discussion"
    sub_para = sub_frame.paragraphs[0]
    sub_para.font.size = Pt(32)
    sub_para.font.color.rgb = COLORS['highlight']
    sub_para.alignment = PP_ALIGN.CENTER

    # Contact Info
    contact_box = slide.shapes.add_textbox(Inches(1.5), Inches(4.8), Inches(7), Inches(1.5))
    contact_frame = contact_box.text_frame

    contact_text = """📧 support@neurovidya.com
🌐 www.neurovidya.com
💻 github.com/neurovidya
📖 docs.neurovidya.com"""

    contact_frame.text = contact_text
    for para in contact_frame.paragraphs:
        para.font.size = Pt(20)
        para.font.color.rgb = COLORS['white']

    # Social
    social_box = slide.shapes.add_textbox(Inches(1), Inches(6.8), Inches(8), Inches(0.5))
    social_frame = social_box.text_frame
    social_frame.text = "🐦 @NeuroVidyaApp | 💬 Discord Community | 📝 blog.neurovidya.com"
    social_para = social_frame.paragraphs[0]
    social_para.font.size = Pt(18)
    social_para.font.color.rgb = COLORS['white']
    social_para.alignment = PP_ALIGN.CENTER

    # Acknowledgments
    ack_box = slide.shapes.add_textbox(Inches(1), Inches(7.5), Inches(8), Inches(0.4))
    ack_frame = ack_box.text_frame
    ack_frame.text = "Special thanks to our mentors, the dyslexia community, and open-source contributors"
    ack_para = ack_frame.paragraphs[0]
    ack_para.font.size = Pt(16)
    ack_para.font.color.rgb = COLORS['white']
    ack_para.alignment = PP_ALIGN.CENTER
    ack_para.font.italic = True

    return slide

def main():
    """Create the complete presentation"""
    print("Creating NeuroVidya Presentation...")

    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Get blank layout
    blank_layout = prs.slide_layouts[6]

    # Add all slides
    print("Adding Slide 1: Title...")
    add_title_slide(prs, blank_layout)

    print("Adding Slide 2: Introduction...")
    add_intro_slide(prs, blank_layout)

    print("Adding Slide 3: Problem Statement...")
    add_problem_slide(prs, blank_layout)

    print("Adding Slide 4: Literature Review...")
    add_lit_review_slide(prs, blank_layout)

    print("Adding Slide 5: Objectives...")
    add_objectives_slide(prs, blank_layout)

    print("Adding Slide 6: Methodology...")
    add_methodology_slide(prs, blank_layout)

    print("Adding Slide 7: Implementation...")
    add_implementation_slide(prs, blank_layout)

    print("Adding Slide 8: Application...")
    add_application_slide(prs, blank_layout)

    print("Adding Slide 9: Conclusion...")
    add_conclusion_slide(prs, blank_layout)

    print("Adding Slide 10: Thank You...")
    add_thankyou_slide(prs, blank_layout)

    # Save presentation
    output_path = "/Users/apple/Desktop/NeuroVidya MIni Project/NeuroVidya_Presentation.pptx"
    prs.save(output_path)

    print(f"\n✅ Presentation created successfully!")
    print(f"📁 Saved to: {output_path}")
    print(f"📊 Total slides: {len(prs.slides)}")
    print("\nDesign Features:")
    print("  • Dyslexia-friendly color scheme (high contrast)")
    print("  • Large, readable fonts (18-54pt)")
    print("  • Clean layouts with ample white space")
    print("  • Sans-serif typography")
    print("  • Icons and visual elements")
    print("  • Left-aligned text for easier reading")

if __name__ == "__main__":
    main()