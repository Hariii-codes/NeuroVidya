from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# Colors
CREAM = RGBColor(0xFA, 0xF6, 0xF0)
DARK_GRAY = RGBColor(0x2D, 0x37, 0x48)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # Standard 16:9
prs.slide_height = Inches(7.5)

# Set default slide layout (blank)
blank_layout = prs.slide_layouts[6]  # Blank layout

def add_textbox(slide, left, top, width, height, text, font_size=24, bold=False, italic=False, color=DARK_GRAY, alignment=PP_ALIGN.LEFT, font_name='Arial'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def set_bg(slide, color=CREAM):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

# Slide 1: Title Slide
slide1 = prs.slides.add_slide(blank_layout)
set_bg(slide1, CREAM)

# Title
add_textbox(slide1, Inches(1), Inches(1.5), Inches(11), Inches(1.5), "NEUROVIDYA", font_size=54, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER, font_name='Arial')
# Subtitle
add_textbox(slide1, Inches(1), Inches(3.2), Inches(11), Inches(1), "A Dyslexia-First Adaptive Learning Platform", font_size=32, bold=False, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
# Tagline
add_textbox(slide1, Inches(1), Inches(4.2), Inches(11), Inches(0.8), "\"Transforming Learning Through Accessibility-First Design\"", font_size=24, italic=True, color=PURPLE, alignment=PP_ALIGN.CENTER)
# Team members (simplified)
add_textbox(slide1, Inches(1), Inches(5.5), Inches(11), Inches(1), "Team Members: [Your Name] - Project Lead, [Member 2] - Backend, [Member 3] - UI/UX, [Member 4] - Research", font_size=18, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Slide 2: Introduction
slide2 = prs.slides.add_slide(blank_layout)
set_bg(slide2, CREAM)
add_textbox(slide2, Inches(0.8), Inches(0.5), Inches(11.733), Inches(1), "Introduction to NeuroVidya", font_size=44, bold=True, color=BLUE)
add_textbox(slide2, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5),
            "• Web-based adaptive learning platform designed for dyslexia\n• Built on dyslexia-first design methodology\n• Integrates AI with evidence-based learning strategies\n• 10-20% of global population affected by dyslexia\n• Accessibility from the ground up, not as an add-on\n• Empowering users through extensive customization\n• Privacy-first approach with browser-native features",
            font_size=24, color=DARK_GRAY)

# Slide 3: Literature Review
slide3 = prs.slides.add_slide(blank_layout)
set_bg(slide3, CREAM)
add_textbox(slide3, Inches(0.8), Inches(0.5), Inches(11.733), Inches(1), "Literature Review & Research Foundation", font_size=44, bold=True, color=BLUE)
add_textbox(slide3, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5),
            "Theoretical Frameworks:\n• Dual Coding Theory (Paivio, 1971)\n• Cognitive Load Theory (Sweller, 1988)\n• Phonological Deficit Hypothesis\n\nExisting Solutions Limitations:\n• Microsoft Immersive Reader - limited customization\n• Learning Ally - audio-focused only\n• Google Read & Write - extension limits\n• Kurzweil 3000 - expensive, dated",
            font_size=22, color=DARK_GRAY)
add_textbox(slide3, Inches(6.5), Inches(1.8), Inches(6), Inches(5),
            "Research Gap:\n• No comprehensive open-source platform\n• Lack of dyslexia-first design\n• Limited AI personalization\n• No browser-native pronunciation analysis\n\nKey Sources:\n• British Dyslexia Association Style Guide\n• \"Good Fonts for Dyslexia\" - Rello & Baeza-Yates\n• WCAG 2.1 Guidelines\n• \"Overcoming Dyslexia\" - Shaywitz",
            font_size=22, color=DARK_GRAY)

# Slide 4: Problem Statement
slide4 = prs.slides.add_slide(blank_layout)
set_bg(slide4, CREAM)
add_textbox(slide4, Inches(0.8), Inches(0.5), Inches(11.733), Inches(1), "Problem Statement", font_size=44, bold=True, color=BLUE)
add_textbox(slide4, Inches(0.8), Inches(1.8), Inches(11.733), Inches(1.2),
            "\"Educational platforms are designed for the neurotypical majority, leaving 10-20% of the population with dyslexia struggling to access quality learning experiences.\"",
            font_size=26, italic=True, color=PURPLE)
add_textbox(slide4, Inches(0.8), Inches(3.2), Inches(11.733), Inches(3.5),
            "Challenges:\n• Visual crowding and pattern interference\n• Phonological processing deficits\n• Working memory limitations\n• Existing platforms treat accessibility as add-ons\n• Expensive proprietary solutions\n• No real-time pronunciation feedback\n\nResearch Question: How can a web-based platform be architected from its foundation to optimize reading comprehension, retention, and engagement for dyslexic learners?",
            font_size=22, color=DARK_GRAY)

# Slide 5: Objectives
slide5 = prs.slides.add_slide(blank_layout)
set_bg(slide5, CREAM)
add_textbox(slide5, Inches(0.8), Inches(0.5), Inches(11.733), Inches(1), "Project Objectives", font_size=44, bold=True, color=BLUE)
add_textbox(slide5, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5),
            "Primary Objectives:\n🎯 Design dyslexia-first design system\n🎯 Develop advanced phonetic support\n🎯 Create multi-modal learning experience\n🎯 Ensure accessibility & usability\n\nSuccess Metrics:\n• WCAG 2.1 AA compliance\n• 90%+ accessibility score\n• <2s page load time\n• 79%+ test coverage",
            font_size=22, color=DARK_GRAY)
add_textbox(slide5, Inches(6.5), Inches(1.8), Inches(6), Inches(5),
            "Secondary Objectives:\n📊 Establish measurable benchmarks\n🔬 Validate through research\n💡 Advance the field\n\nCurrent Status:\n✅ Accessibility: 100%\n✅ Test Coverage: 79%\n✅ Page Load: 2.8s\n✅ Design Components: 20+",
            font_size=22, color=DARK_GRAY)

# Slide 6: Software Requirements
slide6 = prs.slides.add_slide(blank_layout)
set_bg(slide6, CREAM)
add_textbox(slide6, Inches(0.8), Inches(0.5), Inches(11.733), Inches(1), "Software & Hardware Requirements", font_size=44, bold=True, color=BLUE)
add_textbox(slide6, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5),
            "Frontend:\n• React 18.2.0 + TypeScript 5.3.3\n• Vite 5.0.8\n• Modern browser (Chrome, Firefox, Safari)\n\nBackend:\n• Python 3.9+\n• FastAPI 0.104.1\n• PostgreSQL or MongoDB\n\nClient-Side:\n• 4GB RAM (8GB recommended)\n• Microphone for pronunciation\n• Stable internet",
            font_size=22, color=DARK_GRAY)
add_textbox(slide6, Inches(6.5), Inches(1.8), Inches(6), Inches(5),
            "API Services:\n• OpenAI API (GPT-4o-mini)\n• Anthropic API (Claude 3.5 Haiku)\n• Free tiers sufficient\n\nPlatform Support:\n✅ Windows 10/11\n✅ macOS 11+\n✅ Linux\n🔄 Android/iOS (Beta)\n\nAccessibility Tools:\n• JAWS, NVDA, VoiceOver, Narrator",
            font_size=22, color=DARK_GRAY)

# Slide 7: Methodology
slide7 = prs.slides.add_slide(blank_layout)
set_bg(slide7, CREAM)
add_textbox(slide7, Inches(0.8), Inches(0.5), Inches(11.733), Inches(1), "Methodology & Development Approach", font_size=44, bold=True, color=BLUE)
add_textbox(slide7, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5),
            "Agile Development with 2-Week Sprints:\n\n• Phase 1: Research & Planning (Weeks 1-2) ✅\n• Phase 2: Design System Foundation (Weeks 3-4) ✅\n• Phase 3: Advanced Features (Weeks 5-7) ✅\n• Phase 4: Page Development (Weeks 8-10) ✅\n• Phase 5: Backend Integration (Weeks 11-12) ✅\n• Phase 6: Testing & QA (Weeks 13-14) ✅\n• Phase 7: User Validation (Weeks 15-16) 🔄\n\nComponent-First Architecture: Atomic Design Methodology",
            font_size=22, color=DARK_GRAY)

# Slide 8: Implementation
slide8 = prs.slides.add_slide(blank_layout)
set_bg(slide8, CREAM)
add_textbox(slide8, Inches(0.8), Inches(0.5), Inches(11.733), Inches(1), "Implementation Details", font_size=44, bold=True, color=BLUE)
add_textbox(slide8, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5),
            "System Architecture:\n• Client Layer: React SPA + Web Speech API\n• API Layer: FastAPI\n• Data Layer: PostgreSQL + MongoDB\n\nKey Components (20+):\n• DyslexiaProvider, DyslexiaTheme\n• DyslexiaButton, DyslexiaInput\n• ReadingGuide, PhoneticHighlighter\n• Pronunciation analysis\n• Phonetic chunking algorithm",
            font_size=22, color=DARK_GRAY)
add_textbox(slide8, Inches(6.5), Inches(1.8), Inches(6), Inches(5),
            "Tech Stack:\n• Frontend: React, TypeScript, Vite, Zustand\n• Backend: FastAPI, Prisma ORM\n• AI: OpenAI GPT-4o-mini, Claude\n• Testing: Vitest, jest-axe\n\nPerformance Metrics:\n• Bundle: 843 KB (200 KB gzipped)\n• Lighthouse: 92\n• Accessibility: 100%\n• Test Coverage: 79%",
            font_size=22, color=DARK_GRAY)

# Slide 9: Application & Future Scope
slide9 = prs.slides.add_slide(blank_layout)
set_bg(slide9, CREAM)
add_textbox(slide9, Inches(0.8), Inches(0.5), Inches(11.733), Inches(1), "Application & Future Scope", font_size=44, bold=True, color=BLUE)
add_textbox(slide9, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5),
            "Current Applications:\n🎓 K-12 Schools - Special education\n🎓 Higher Education - Disability services\n💼 Workplace accommodation\n👤 Individual learners\n\nFor Students:\n✅ Personalized reading experience\n✅ Real-time pronunciation feedback\n✅ Interactive learning games\n✅ Progress tracking",
            font_size=22, color=DARK_GRAY)
add_textbox(slide9, Inches(6.5), Inches(1.8), Inches(6), Inches(5),
            "Future Enhancements:\n• Mobile Apps (React Native)\n• Advanced Analytics\n• Multi-language Support\n• AI Personalized Learning Paths\n• LMS Integrations (Canvas, Blackboard)\n\nGlobal Impact:\n• 700+ million people with dyslexia\n• Scalable open-source solution\n• Improved academic outcomes",
            font_size=22, color=DARK_GRAY)

# Slide 10: Conclusion & Thank You
slide10 = prs.slides.add_slide(blank_layout)
set_bg(slide10, CREAM)
add_textbox(slide10, Inches(1), Inches(1), Inches(11), Inches(1.2), "Conclusion & Thank You", font_size=48, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)
add_textbox(slide10, Inches(1), Inches(2.5), Inches(11), Inches(1),
            "NeuroVidya: A paradigm shift from accessibility-as-accommodation to accessibility-as-foundation.",
            font_size=26, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)
add_textbox(slide10, Inches(1), Inches(4), Inches(11), Inches(1),
            "🎯 Technical Excellence ✨ Pedagogical Innovation 🎯 User Empowerment",
            font_size=24, color=PURPLE, alignment=PP_ALIGN.CENTER)
add_textbox(slide10, Inches(1), Inches(5.5), Inches(11), Inches(1),
            "Contact: support@neurovidya.com | www.neurovidya.com | github.com/neurovidya",
            font_size=20, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# Save presentation
output_path = '/Users/apple/Desktop/NeuroVidya MIni Project/NeuroVidya_Presentation.pptx'
prs.save(output_path)
print(f"Presentation saved to {output_path}")
